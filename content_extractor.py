import csv
import io
import ipaddress
import logging
import re
import socket
from pathlib import Path
from typing import Optional
from urllib.parse import urlparse

import requests
from bs4 import BeautifulSoup
from PIL import Image

logger = logging.getLogger("rcjy.content_extractor")

# SSRF limits
_ALLOWED_SCHEMES = {"http", "https"}
_MAX_URL_RESPONSE_BYTES = 10 * 1024 * 1024  # 10 MB max download from URL
_URL_REQUEST_TIMEOUT = 15  # seconds
_MAX_FILE_SIZE = 50 * 1024 * 1024  # 50 MB per uploaded file


_BLOCKED_HOSTS = {
    "metadata.google.internal",
    "metadata.google",
    "169.254.169.254",
    "100.100.100.200",
}


def _is_safe_ip(ip_str: str) -> bool:
    # Check if IP is safe (not private/reserved/loopback)
    try:
        ip = ipaddress.ip_address(ip_str)
        return not (
            ip.is_private
            or ip.is_loopback
            or ip.is_link_local
            or ip.is_reserved
            or ip.is_multicast
            or ip.is_unspecified
        )
    except ValueError:
        return False


def _resolve_and_validate(hostname: str) -> str:
    # Resolve hostname once and return a safe IP, or raise ValueError
    if hostname.lower() in _BLOCKED_HOSTS:
        raise ValueError("Access to cloud metadata endpoints is not allowed.")
    try:
        infos = socket.getaddrinfo(hostname, None, socket.AF_UNSPEC, socket.SOCK_STREAM)
    except (socket.gaierror, OSError):
        raise ValueError("Could not resolve hostname.")
    for info in infos:
        ip_str = info[4][0]
        if _is_safe_ip(ip_str):
            return ip_str
    raise ValueError("URLs pointing to internal/private network addresses are not allowed.")


def _validate_url(url: str) -> tuple[str, str]:
    # Validate URL, resolve DNS once. Returns (url, resolved_ip)
    url = url.strip()
    if not url:
        raise ValueError("Empty URL")
    parsed = urlparse(url)
    if parsed.scheme.lower() not in _ALLOWED_SCHEMES:
        raise ValueError(f"URL scheme '{parsed.scheme}' is not allowed. Only http/https are permitted.")
    hostname = parsed.hostname
    if not hostname:
        raise ValueError("URL has no hostname.")
    resolved_ip = _resolve_and_validate(hostname)
    return url, resolved_ip

try:
    from pypdf import PdfReader
except ImportError:
    try:
        from PyPDF2 import PdfReader
    except ImportError:
        PdfReader = None

try:
    from docx import Document as DocxDocument
except ImportError:
    DocxDocument = None

MIME_MAP = {
    ".jpg": "image/jpeg", ".jpeg": "image/jpeg", ".png": "image/png",
    ".gif": "image/gif", ".webp": "image/webp", ".svg": "image/svg+xml",
    ".bmp": "image/bmp", ".tiff": "image/tiff", ".tif": "image/tiff",
    ".ico": "image/x-icon",
    ".mp3": "audio/mpeg", ".wav": "audio/wav", ".ogg": "audio/ogg",
    ".m4a": "audio/mp4", ".flac": "audio/flac", ".aac": "audio/aac",
    ".wma": "audio/x-ms-wma",
    ".mp4": "video/mp4", ".avi": "video/x-msvideo", ".mov": "video/quicktime",
    ".mkv": "video/x-matroska", ".webm": "video/webm", ".wmv": "video/x-ms-wmv",
    ".flv": "video/x-flv", ".mpeg": "video/mpeg",
    ".pdf": "application/pdf",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".doc": "application/msword",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ".ppt": "application/vnd.ms-powerpoint",
    ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ".xls": "application/vnd.ms-excel",
    ".csv": "text/csv", ".json": "application/json", ".xml": "application/xml",
    ".html": "text/html", ".htm": "text/html",
    ".txt": "text/plain", ".md": "text/markdown", ".rtf": "application/rtf",
}


def get_mime_type(filename: str) -> str:
    return MIME_MAP.get(Path(filename).suffix.lower(), "application/octet-stream")


def is_image(filename: str) -> bool:
    return get_mime_type(filename).startswith("image/")


def is_audio(filename: str) -> bool:
    return get_mime_type(filename).startswith("audio/")


def is_video(filename: str) -> bool:
    return get_mime_type(filename).startswith("video/")


def extract_from_url(url: str, max_chars: int = 50000) -> str:
    try:
        url, resolved_ip = _validate_url(url)
    except ValueError as e:
        logger.warning("URL validation failed for user-supplied URL: %s", e)
        return f"Invalid URL: {e}"

    try:
        # Pin the resolved IP to prevent DNS rebinding (TOCTOU)
        parsed = urlparse(url)
        hostname = parsed.hostname
        port = parsed.port or (443 if parsed.scheme == "https" else 80)
        pinned_url = url.replace(f"://{hostname}", f"://{resolved_ip}", 1)

        headers = {
            "User-Agent": "Mozilla/5.0 (compatible; RCJY-MediaBot/1.0)",
            "Host": hostname,
        }
        session = requests.Session()
        session.max_redirects = 5
        # Disable redirects to validate each hop
        response = session.get(
            pinned_url,
            headers=headers,
            timeout=_URL_REQUEST_TIMEOUT,
            allow_redirects=False,
            stream=True,
            verify=True if parsed.scheme == "https" else False,
        )

        # Follow redirects manually with validation
        redirect_count = 0
        while response.is_redirect and redirect_count < 5:
            redirect_count += 1
            next_url = response.headers.get("Location", "")
            if not next_url:
                break
            # Resolve absolute URL
            if not next_url.startswith("http"):
                next_url = f"{parsed.scheme}://{hostname}{next_url}"
            try:
                next_url, next_ip = _validate_url(next_url)
            except ValueError:
                logger.warning("Redirect to blocked destination: %s", next_url)
                return "Error: URL redirected to a blocked destination."
            next_parsed = urlparse(next_url)
            next_host = next_parsed.hostname
            pinned_next = next_url.replace(f"://{next_host}", f"://{next_ip}", 1)
            response = session.get(
                pinned_next,
                headers={"User-Agent": headers["User-Agent"], "Host": next_host},
                timeout=_URL_REQUEST_TIMEOUT,
                allow_redirects=False,
                stream=True,
            )

        response.raise_for_status()

        # Check Content-Length
        content_length = response.headers.get("Content-Length")
        if content_length and int(content_length) > _MAX_URL_RESPONSE_BYTES:
            return "Error: URL response exceeds maximum allowed size (10 MB)."

        # Read with size limit
        chunks = []
        total = 0
        for chunk in response.iter_content(chunk_size=8192, decode_unicode=False):
            total += len(chunk)
            if total > _MAX_URL_RESPONSE_BYTES:
                logger.warning("URL response exceeded size limit, truncating")
                break
            chunks.append(chunk)
        raw_content = b"".join(chunks)

        text_content = raw_content.decode("utf-8", errors="ignore")
        soup = BeautifulSoup(text_content, "html.parser")
        for tag in soup(["script", "style", "nav", "footer", "header"]):
            tag.decompose()
        text = soup.get_text(separator="\n", strip=True)
        text = re.sub(r"\n{3,}", "\n\n", text)
        if len(text) > max_chars:
            text = text[:max_chars] + "\n\n[Content truncated...]"
        return text.strip() or "Could not extract text from URL."
    except requests.exceptions.Timeout:
        return "Error: URL request timed out."
    except requests.exceptions.ConnectionError:
        return "Error: Could not connect to URL."
    except requests.exceptions.HTTPError as e:
        status = e.response.status_code if e.response is not None else "unknown"
        return f"Error: URL returned HTTP {status}."
    except Exception:
        logger.exception("Unexpected error fetching URL")
        return "Error: Could not fetch content from URL."


def extract_from_pdf(file) -> str:
    if PdfReader is None:
        return "Error: PDF reader not available."
    try:
        reader = PdfReader(file)
        parts = [page.extract_text() for page in reader.pages if page.extract_text()]
        return "\n\n".join(parts)[:50000] or "No text found in PDF."
    except Exception:
        logger.exception("Error reading PDF")
        return "Error: Could not read PDF file."


def extract_from_docx(file) -> str:
    if DocxDocument is None:
        return "Error: DOCX reader not available."
    try:
        doc = DocxDocument(file)
        return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())[:50000] or "No text found."
    except Exception:
        logger.exception("Error reading DOCX")
        return "Error: Could not read DOCX file."


def extract_from_txt(file) -> str:
    try:
        content = file.read()
        if isinstance(content, bytes):
            content = content.decode("utf-8", errors="ignore")
        return content[:50000]
    except Exception:
        logger.exception("Error reading text file")
        return "Error: Could not read text file."


def extract_from_csv(file) -> str:
    try:
        content = file.read()
        if isinstance(content, bytes):
            content = content.decode("utf-8", errors="ignore")
        reader = csv.reader(io.StringIO(content))
        rows = []
        for i, row in enumerate(reader):
            if i > 500:
                rows.append("[... truncated ...]")
                break
            rows.append(" | ".join(row))
        return "\n".join(rows)[:50000]
    except Exception:
        logger.exception("Error reading CSV")
        return "Error: Could not read CSV file."


def extract_from_pptx(file) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(file)
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            slide_text.append(para.text.strip())
            if slide_text:
                parts.append(f"[Slide {i}]\n" + "\n".join(slide_text))
        return "\n\n".join(parts)[:50000] or "No text found in presentation."
    except ImportError:
        return "Error: PPTX reader not available."
    except Exception:
        logger.exception("Error reading PPTX")
        return "Error: Could not read PPTX file."


def extract_from_xlsx(file) -> str:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(file, data_only=True)
        parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(max_row=200, values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(cells):
                    rows.append(" | ".join(cells))
            if rows:
                parts.append(f"[Sheet: {sheet_name}]\n" + "\n".join(rows))
        return "\n\n".join(parts)[:50000] or "No data found in spreadsheet."
    except ImportError:
        return "Error: Excel reader not available."
    except Exception:
        logger.exception("Error reading Excel")
        return "Error: Could not read Excel file."


def get_content_from_input(
    text: Optional[str] = None,
    url: Optional[str] = None,
    files: Optional[list] = None,
) -> tuple[str, list[tuple[str, str, bytes]]]:
    parts: list[tuple[str, str]] = []
    attachments: list[tuple[str, str, bytes]] = []

    if text and text.strip():
        parts.append(("User input", text.strip()))

    if url and url.strip():
        parts.append(("URL content", extract_from_url(url.strip())))

    if files:
        for f in files:
            if f is None:
                continue
            name = getattr(f, "name", "attachment")
            # Strip path components to prevent traversal
            name = Path(name).name  # removes any directory components
            suffix = Path(name).suffix.lower()
            mime = get_mime_type(name)

            try:
                f.seek(0)
            except Exception:
                pass

            # Check file size
            try:
                f.seek(0, 2)  # seek to end
                file_size = f.tell()
                f.seek(0)
                if file_size > _MAX_FILE_SIZE:
                    logger.warning("File '%s' exceeds size limit (%d bytes)", name, file_size)
                    parts.append((f"File: {name}", f"[File too large: {file_size:,} bytes, limit is {_MAX_FILE_SIZE:,}]"))
                    continue
            except Exception:
                pass  # If we can't check size, proceed cautiously

            if is_image(name):
                try:
                    data = f.read()
                    attachments.append((name, mime, data))
                    parts.append((f"Image: {name}", "[Image attached]"))
                except Exception:
                    pass

            elif is_audio(name):
                try:
                    data = f.read()
                    attachments.append((name, mime, data))
                    parts.append((f"Audio: {name}", f"[Audio file attached: {name}]"))
                except Exception:
                    pass

            elif is_video(name):
                try:
                    data = f.read()
                    attachments.append((name, mime, data))
                    parts.append((f"Video: {name}", f"[Video file attached: {name}]"))
                except Exception:
                    pass

            elif suffix == ".pdf":
                content = extract_from_pdf(f)
                parts.append((f"PDF: {name}", content))
                try:
                    f.seek(0)
                    attachments.append((name, mime, f.read()))
                except Exception:
                    pass

            elif suffix in (".docx", ".doc"):
                content = extract_from_docx(f)
                parts.append((f"Document: {name}", content))
                try:
                    f.seek(0)
                    attachments.append((name, mime, f.read()))
                except Exception:
                    pass

            elif suffix in (".pptx", ".ppt"):
                content = extract_from_pptx(f)
                parts.append((f"Presentation: {name}", content))
                try:
                    f.seek(0)
                    attachments.append((name, mime, f.read()))
                except Exception:
                    pass

            elif suffix in (".xlsx", ".xls"):
                content = extract_from_xlsx(f)
                parts.append((f"Spreadsheet: {name}", content))
                try:
                    f.seek(0)
                    attachments.append((name, mime, f.read()))
                except Exception:
                    pass

            elif suffix == ".csv":
                parts.append((f"CSV: {name}", extract_from_csv(f)))

            elif suffix in (".txt", ".md", ".rtf", ".json", ".xml", ".html", ".htm"):
                parts.append((f"Text: {name}", extract_from_txt(f)))

            else:
                try:
                    raw = f.read()
                    if isinstance(raw, bytes):
                        if len(raw) < 200_000:
                            try:
                                txt = raw.decode("utf-8", errors="strict")
                                parts.append((f"File: {name}", txt[:30000]))
                            except UnicodeDecodeError:
                                attachments.append((name, mime, raw))
                                parts.append((f"Binary: {name}", f"[Binary file attached: {name}]"))
                        else:
                            attachments.append((name, mime, raw))
                            parts.append((f"File: {name}", f"[Large file attached: {name}]"))
                except Exception:
                    pass

    combined = "\n\n---\n\n".join(f"[{title}]\n{content}" for title, content in parts)
    return combined or "No content provided.", attachments
